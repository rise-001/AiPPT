from typing import Dict
import json
from textwrap import dedent
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
import re
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
from gemini_genai import gen_image, gen_json_text, gen_text
def gen_outline(idea_prompt:str)->list[dict]:
    """generate outline of ppt, including optional parts and pages with title and points"""
    outline_prompt = dedent(f"""\
    You are a helpful assistant that generates an outline for a ppt.
    
    You can organize the content in two ways:
    
    1. Simple format (for short PPTs without major sections):
    [{{"title": "title1", "points": ["point1", "point2"]}}, {{"title": "title2", "points": ["point1", "point2"]}}]
    
    2. Part-based format (for longer PPTs with major sections):
    [
      {{
        "part": "Part 1: Introduction",
        "pages": [
          {{"title": "Welcome", "points": ["point1", "point2"]}},
          {{"title": "Overview", "points": ["point1", "point2"]}}
        ]
      }},
      {{
        "part": "Part 2: Main Content",
        "pages": [
          {{"title": "Topic 1", "points": ["point1", "point2"]}},
          {{"title": "Topic 2", "points": ["point1", "point2"]}}
        ]
      }}
    ]
    
    Choose the format that best fits the content. Use parts when the PPT has clear major sections.
    
    The user's request: {idea_prompt}. Now generate the outline, don't include any other text.
    使用全中文输出。
    """)
    outline = gen_json_text(outline_prompt)
    outline = json.loads(outline)
    return outline
    

def flatten_outline(outline: list[dict]) -> list[dict]:
    """将可能包含part结构的outline扁平化为页面列表"""
    pages = []
    for item in outline:
        if "part" in item and "pages" in item:
            # 这是一个part，展开其中的页面
            for page in item["pages"]:
                # 为每个页面添加part信息
                page_with_part = page.copy()
                page_with_part["part"] = item["part"]
                pages.append(page_with_part)
        else:
            # 这是一个直接的页面
            pages.append(item)
    return pages

def gen_desc(idea_prompt, outline: list[Dict])->list[Dict] :
    """generate description for each page, including title, full text content and more (并行生成)"""
    # 先将outline扁平化为页面列表
    pages = flatten_outline(outline)
    
    # 为每个页面准备生成任务
    def generate_page_desc(i, page_outline):
        part_info = f"\nThis page belongs to: {page_outline['part']}" if 'part' in page_outline else ""
        desc_prompt = dedent(f"""\
        we are generating the text desciption for each ppt page.
        the original user request is: \n{idea_prompt}\n
        We already have the entire ouline: \n{outline}\n{part_info}
        Now please generate the description for page {i}:
        {page_outline}
        The description includes page title, text to render(keep it concise).
        For example:
        页面标题：原始社会：与自然共生
        页面文字：
        - 狩猎采集文明： 人类活动规模小，对环境影响有限。
        - 依赖性强： 生活完全依赖于自然资源的直接供给，对自然规律敬畏。
        - 适应而非改造： 通过观察和模仿学习自然，发展出适应当地环境的生存技能。
        - 影响特点： 局部、短期、低强度，生态系统有充足的自我恢复能力。
        
        使用全中文输出。
        """)
        page_desc = gen_text(desc_prompt)
        # 清理多余的缩进
        page_desc = dedent(page_desc)
        return (i, page_desc)  # 返回索引和描述，以便排序
    
    # 使用线程池并行生成所有页面的描述
    desc_dict = {}
    with ThreadPoolExecutor(max_workers=5) as executor:
        # 提交所有任务
        futures = [executor.submit(generate_page_desc, i, page_outline) 
                   for i, page_outline in enumerate(pages, 1)]
        
        # 收集结果
        for future in as_completed(futures):
            i, page_desc = future.result()
            desc_dict[i] = page_desc
            print(f"✓ 页面 {i}/{len(pages)} 描述生成完成")
    
    # 按照原始顺序返回结果
    desc = [desc_dict[i] for i in sorted(desc_dict.keys())]
    return desc

def gen_outline_text(outline: list[Dict]) -> str:
    """将outline转换为文本格式，用于提示词"""
    text_parts = []
    for i, item in enumerate(outline, 1):
        if "part" in item and "pages" in item:
            text_parts.append(f"{i}. {item['part']}")
        else:
            text_parts.append(f"{i}. {item.get('title', 'Untitled')}")
    result = "\n".join(text_parts)
    # 清理多余的缩进
    return dedent(result)

def gen_prompts(outline: list[Dict], desc: list[str]) -> list[str]:
    """为每页描述生成图片提示词"""
    pages = flatten_outline(outline)
    outline_text = gen_outline_text(outline)
    
    prompts = []
    for i, (page, page_desc) in enumerate(zip(pages, desc), 1):
        # 确定当前所属章节
        if 'part' in page:
            current_section = page['part']
        else:
            current_section = f"{page.get('title', 'Untitled')}"
        
        # 构建提示词，参考generate-example.py的格式
        prompt = dedent(f"""\
        利用专业平面设计知识，根据参考图片的色彩与风格生成一页设计风格相同的ppt页面，作为整个ppt的其中一页，内容是:
        {page_desc}
        
        整个ppt的大纲为：
        {outline_text}
        
        当前位于章节：{current_section}
        
        要求文字清晰锐利，画面为4k分辨率 16:9比例.画面风格与配色保持严格一致。ppt使用全中文。
        """)
        print(f"\n-----\n prompt{i}:\n {prompt}\n-----\n")
        prompts.append(prompt)
    
    return prompts

def gen_images_parallel(prompts: list[str], ref_image: str, output_dir: str = "output") -> list[str]:
    """并行生成所有PPT页面图片"""
    # 创建输出目录
    output_path = Path(output_dir)
    output_path.mkdir(exist_ok=True)
    
    def generate_single_image(i, prompt):
        """生成单张图片"""
        try:
            print(f"🎨 开始生成页面 {i}/{len(prompts)} 的图片...")
            image = gen_image(prompt, ref_image)
            if image:
                output_file = output_path / f"slide_{i:02d}.png"
                image.save(str(output_file))
                print(f"✓ 页面 {i}/{len(prompts)} 图片生成完成: {output_file}")
                return (i, str(output_file))
            else:
                print(f"✗ 页面 {i}/{len(prompts)} 图片生成失败")
                return (i, None)
        except Exception as e:
            print(f"✗ 页面 {i}/{len(prompts)} 生成出错: {e}")
            return (i, None)
    
    # 使用线程池并行生成所有图片
    image_files = {}
    with ThreadPoolExecutor(max_workers=8) as executor:  # 限制并发数为3避免API限流
        # 提交所有任务
        futures = [executor.submit(generate_single_image, i, prompt) 
                   for i, prompt in enumerate(prompts, 1)]
        
        # 收集结果
        for future in as_completed(futures):
            i, image_file = future.result()
            image_files[i] = image_file
    
    # 按照原始顺序返回结果
    return [image_files[i] for i in sorted(image_files.keys())]

def create_pptx_from_images(input_dir: str = "output", output_file: str = "presentation.pptx"):
    """
    将指定目录下的slide_XX.png图片按顺序组合成PPTX文件
    
    Args:
        input_dir: 输入图片所在目录
        output_file: 输出的PPTX文件名
    """
    input_path = Path(input_dir)
    slide_files = list(input_path.glob("slide_*.png"))
    
    def extract_number(filename):
        match = re.search(r'slide_(\d+)', filename.stem)
        return int(match.group(1)) if match else 0
    
    slide_files.sort(key=extract_number)
    
    print(f"\n📁 找到 {len(slide_files)} 张幻灯片图片")
    print(f"📝 开始创建 PPTX 文件...")
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸为16:9 (宽10英寸，高5.625英寸)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 为每张图片创建一页幻灯片
    for i, image_file in enumerate(slide_files, 1):
        print(f"  ✓ 添加第 {i} 页: {image_file.name}")
        
        # 添加空白幻灯片布局（完全空白，没有任何占位符）
        blank_slide_layout = prs.slide_layouts[6]  # 布局6通常是空白布局
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 将图片添加到幻灯片，填充整个页面
        # 左上角位置(0,0)，尺寸为幻灯片的完整宽高
        slide.shapes.add_picture(
            str(image_file),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height
        )
    
    # 保存PPTX文件
    prs.save(output_file)
    
    print(f"\n✅ 成功创建 PPTX 文件: {output_file}")
    print(f"📊 总共 {len(slide_files)} 页幻灯片")
    return True

def gen_ppt(idea_prompt, ref_image):
    # 创建带时间戳的输出目录
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_dir = f"output_{timestamp}"
    pptx_filename = f"presentation_{timestamp}.pptx"
    
    print(f"📂 本次运行输出目录: {output_dir}")
    print(f"📄 PPTX文件名: {pptx_filename}\n")
    
    outline = gen_outline(idea_prompt)
    
    # 显示原始outline结构（可能包含parts）
    print("PPT Outline:")
    for item in outline:
        if "part" in item and "pages" in item:
            print(f"\n【{item['part']}】")
            for j, page in enumerate(item["pages"], 1):
                print(f"  Page {j}: {page.get('title', 'Untitled')}")
                print(f"    Points: {page.get('points', [])}")
        else:
            print(f"\nPage: {item.get('title', 'Untitled')}")
            print(f"  Points: {item.get('points', [])}")
    
    # 生成详细描述
    desc = gen_desc(idea_prompt, outline)
    
    # 显示每页描述
    pages = flatten_outline(outline)
    for i, (page, page_desc) in enumerate(zip(pages, desc), 1):
        part_tag = f"[{page['part']}] " if 'part' in page else ""
        print(f"-----\nPage {i} {part_tag}- {page.get('title', 'Untitled')}\n-----")
        print(f"{page_desc}\n")
    
    # 生成图片提示词
    print("开始生成图片提示词...")
    prompts = gen_prompts(outline, desc)
    print(f"✓ 已生成 {len(prompts)} 个页面的提示词\n")
    
    # 并行生成所有页面图片（使用带时间戳的目录）
    print("开始并行生成PPT页面图片...")
    image_files = gen_images_parallel(prompts, ref_image, output_dir)
    
    # 显示结果汇总
    print("PPT图片生成完成！")
    successful = [f for f in image_files if f is not None]
    print(f"✓ 成功生成 {len(successful)}/{len(image_files)} 张图片")
    for i, image_file in enumerate(image_files, 1):
        if image_file:
            print(f"  页面 {i}: {image_file}")
        else:
            print(f"  页面 {i}: 生成失败")
    
    # 将所有图片组合成PPTX文件
    if successful:
        print("正在生成最终的PPTX文件...")
        create_pptx_from_images(output_dir, pptx_filename)
    
    return image_files
    
    

if __name__ == "__main__":
    idea_prompt="生成一张关于人类活动对生态环境影响的ppt.只要3页。"
    ref_image="template_g.png"
    gen_ppt(idea_prompt, ref_image)